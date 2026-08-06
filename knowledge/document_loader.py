"""Document loader: lightweight file ingestion for supported formats."""
from __future__ import annotations

import csv
import email
import html
import os
import re
from pathlib import Path
from typing import Optional


class DocumentLoadError(Exception):
    """Raised when a document cannot be loaded."""


def load_document(path: str | Path, *, encoding: str = "utf-8", max_size: int = 20 * 1024 * 1024) -> tuple[str, dict]:
    """Return document body text and basic metadata."""
    p = Path(path)
    if not p.exists() or not p.is_file():
        raise DocumentLoadError(f"Missing document: {p}")

    if p.stat().st_size > max_size:
        raise DocumentLoadError(f"Document too large: {p}")

    suffix = p.suffix.lower()
    text = ""
    metadata = {
        "source": str(p.resolve()),
        "filename": p.name,
        "extension": suffix,
        "size": p.stat().st_size,
        "created": p.stat().st_ctime,
        "modified": p.stat().st_mtime,
        "language": _guess_language(p),
    }

    try:
        if suffix == ".txt":
            text = p.read_text(encoding=encoding, errors="replace")
        elif suffix == ".md":
            text = p.read_text(encoding=encoding, errors="replace")
        elif suffix == ".html":
            text = _html_to_text(p.read_bytes())
            metadata["mime"] = "text/html"
        elif suffix == ".json":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "application/json"
        elif suffix == ".csv":
            text = _csv_to_text(p, encoding=encoding)
            metadata["mime"] = "text/csv"
        elif suffix == ".py":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/x-python"
        elif suffix == ".java":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/x-java"
        elif suffix == ".js":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/javascript"
        elif suffix == ".c":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/x-c"
        elif suffix == ".cpp":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/x-c++"
        elif suffix == ".h":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/x-c"
        elif suffix == ".log":
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/plain"
        elif suffix == ".rtf":
            text = _rtf_to_text(p.read_bytes())
            metadata["mime"] = "text/rtf"
        elif suffix == ".pdf":
            text = _pdf_to_text(p)
            metadata["mime"] = "application/pdf"
        elif suffix == ".docx":
            text = _docx_to_text(p)
            metadata["mime"] = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif suffix in {".xlsx", ".xls"}:
            text = _excel_to_text(p, suffix)
            metadata["mime"] = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" if suffix == ".xlsx" else "application/vnd.ms-excel"
        elif suffix in {".pptx", ".ppt"}:
            text = _ppt_to_text(p, suffix)
            metadata["mime"] = "application/vnd.openxmlformats-officedocument.presentationml.presentation" if suffix == ".pptx" else "application/vnd.ms-powerpoint"
        elif suffix == ".eml":
            text = _eml_to_text(p.read_bytes())
            metadata["mime"] = "message/rfc822"
        else:
            text = p.read_text(encoding=encoding, errors="replace")
            metadata["mime"] = "text/plain"
    except DocumentLoadError:
        raise
    except Exception as exc:
        raise DocumentLoadError(f"Failed to load document: {p}") from exc

    text = _normalize_whitespace(text)
    metadata.setdefault("char_count", len(text))
    metadata.setdefault("title", _guess_title(text, metadata.get("filename", "")))
    return text, metadata


def _normalize_whitespace(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _guess_language(p: Path) -> str:
    try:
        from langdetect import detect  # type: ignore
    except Exception:
        return "unknown"
    sample = ""
    try:
        sample = p.read_text(encoding="utf-8", errors="replace")[:2000]
    except Exception:
        return "unknown"
    if not sample.strip():
        return "unknown"
    try:
        return detect(sample)
    except Exception:
        return "unknown"


def _guess_title(text: str, filename: str) -> str:
    for line in text.splitlines()[:20]:
        line = line.strip()
        if line.startswith("#"):
            return re.sub(r"^#+\s*", "", line).strip()
        if len(line) > 10 and len(line) < 200:
            return line
    return filename


def _html_to_text(data: bytes) -> str:
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(data or b"", "html.parser")
        return soup.get_text(separator="\n")
    except Exception:
        return ""


def _csv_to_text(path: Path, encoding: str = "utf-8") -> str:
    out: list[str] = []
    with path.open("r", encoding=encoding, errors="replace", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            out.append(" | ".join(row))
    return "\n".join(out)


def _pdf_to_text(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
        reader = PdfReader(str(path))
        parts: list[str] = []
        for page in reader.pages:
            try:
                txt = page.extract_text() or ""
                parts.append(txt)
            except Exception:
                pass
        return "\n".join(parts)
    except Exception as exc:
        raise DocumentLoadError(f"PDF load failed: {path}") from exc


def _docx_to_text(path: Path) -> str:
    try:
        from docx import Document  # type: ignore
        doc = Document(str(path))
        lines: list[str] = []
        for para in doc.paragraphs:
            if para.text:
                lines.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception as exc:
        raise DocumentLoadError(f"DOCX load failed: {path}") from exc


def _excel_to_text(path: Path, suffix: str) -> str:
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        out: list[str] = []
        for ws in wb.worksheets:
            out.append(f"[{ws.title}]")
            for row in ws.iter_rows(values_only=True):
                vals = ["" if v is None else str(v) for v in row]
                out.append(" | ".join(vals))
        return "\n".join(out)
    except Exception:
        try:
            import xlrd  # type: ignore
            book = xlrd.open_workbook(str(path))
            out: list[str] = []
            for ws in book.sheets():
                out.append(f"[{ws.name}]")
                for rx in range(ws.nrows):
                    row = [str(ws.cell_value(rx, cx)) for cx in range(ws.ncols)]
                    out.append(" | ".join(row))
            return "\n".join(out)
        except Exception as exc:
            raise DocumentLoadError(f"Excel load failed: {path}") from exc


def _ppt_to_text(path: Path, suffix: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        out: list[str] = []
        for slide in prs.slides:
            out.append("[Slide]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(shape.text)
        return "\n".join(out)
    except Exception as exc:
        raise DocumentLoadError(f"PPT load failed: {path}") from exc


def _rtf_to_text(data: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
        return rtf_to_text(data.decode("utf-8", errors="replace"))
    except Exception:
        try:
            text = data.decode("utf-8", errors="replace")
            text = re.sub(r"\\[a-zA-Z]+[0-9-]*\s?", " ", text)
            text = re.sub(r"[{}]", "", text)
            return html.unescape(text)
        except Exception:
            return ""


def _eml_to_text(data: bytes) -> str:
    try:
        msg = email.message_from_bytes(data)
        parts: list[str] = []
        if msg.get("subject"):
            parts.append(f"Subject: {msg['subject']}")
        if msg.get("from"):
            parts.append(f"From: {msg['from']}")
        if msg.get("to"):
            parts.append(f"To: {msg['to']}")
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_payload(decode=True).decode("utf-8", errors="replace")
                    break
        else:
            body = msg.get_payload(decode=True).decode("utf-8", errors="replace")
        parts.append(body)
        return "\n".join(parts)
    except Exception:
        return ""
